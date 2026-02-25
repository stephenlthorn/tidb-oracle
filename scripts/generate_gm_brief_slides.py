#!/usr/bin/env python3
"""Generate a Google Slides-compatible GM brief deck as a PPTX file."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUTPUT_DIR = ROOT / "presentations"
OUTPUT_PATH = OUTPUT_DIR / "tidb-oracle-gm-brief.pptx"
LOGO_PATH = ROOT / "ui" / "public" / "tidb-logo.png"


COLORS = {
    "bg": RGBColor(246, 248, 251),
    "panel": RGBColor(255, 255, 255),
    "text": RGBColor(24, 32, 45),
    "muted": RGBColor(92, 105, 123),
    "border": RGBColor(220, 226, 236),
    "brand": RGBColor(233, 48, 57),
    "brand_dark": RGBColor(24, 62, 150),
    "ok": RGBColor(22, 156, 82),
}


def style_background(slide) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = COLORS["bg"]


def add_header(slide, title: str, subtitle: str = "") -> None:
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(0.92))
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(255, 255, 255)
    bar.line.color.rgb = COLORS["border"]

    if LOGO_PATH.exists():
        slide.shapes.add_picture(str(LOGO_PATH), Inches(0.36), Inches(0.16), height=Inches(0.58))

    title_box = slide.shapes.add_textbox(Inches(1.22), Inches(0.14), Inches(8.8), Inches(0.45))
    p = title_box.text_frame.paragraphs[0]
    p.text = title
    p.font.name = "Aptos"
    p.font.size = Pt(25)
    p.font.bold = True
    p.font.color.rgb = COLORS["brand_dark"]

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(1.25), Inches(0.56), Inches(10), Inches(0.26))
        p2 = sub_box.text_frame.paragraphs[0]
        p2.text = subtitle
        p2.font.name = "Aptos"
        p2.font.size = Pt(12)
        p2.font.color.rgb = COLORS["muted"]


def add_panel(slide, left: float, top: float, width: float, height: float):
    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    panel.fill.solid()
    panel.fill.fore_color.rgb = COLORS["panel"]
    panel.line.color.rgb = COLORS["border"]
    panel.line.width = Pt(1.0)
    return panel


def panel_title(slide, text: str, left: float, top: float, width: float):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(0.35))
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.font.name = "Aptos"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLORS["text"]


def add_bullets(slide, bullets: list[str], left: float, top: float, width: float, height: float, size: int = 16):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.name = "Aptos"
        p.font.size = Pt(size)
        p.font.color.rgb = COLORS["text"]
        p.space_after = Pt(8)


def slide_1_overview(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    style_background(slide)
    add_header(slide, "TiDB Oracle", "Internal GTM copilot for answers, recommendations, and asset generation")

    add_panel(slide, 0.7, 1.3, 12.0, 5.4)
    panel_title(slide, "What it is and why it matters", 1.0, 1.7, 11.4)
    add_bullets(
        slide,
        [
            "One internal workspace that turns docs + call transcripts into grounded answers and next-step actions.",
            "Outputs include follow-up emails, discovery questions, technical collateral, and account briefs.",
            "Value to GTM: faster deal cycles, consistent technical messaging, and less manual prep per opportunity.",
            "Every answer is evidence-backed with citations and full audit logs.",
        ],
        1.0,
        2.1,
        11.4,
        4.5,
        16,
    )


def slide_2_single_llm_mcp(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    style_background(slide)
    add_header(slide, "Simple Architecture: One Enterprise LLM + MCP", "Single-model path with internal TiDB knowledge grounding")

    add_panel(slide, 0.7, 1.4, 5.85, 4.85)
    panel_title(slide, "Single LLM path", 1.0, 1.78, 5.2)
    add_bullets(
        slide,
        [
            "Use one enterprise-grade OpenAI model endpoint for all generation.",
            "Business/API data is not used for model training by default.",
            "Apply retention and storage controls from enterprise/API policy.",
            "No consumer ChatGPT plan is used for this workload.",
        ],
        1.0,
        2.15,
        5.2,
        3.8,
        14,
    )

    add_panel(slide, 6.85, 1.4, 5.85, 4.85)
    panel_title(slide, "MCP internal knowledge layer", 7.15, 1.78, 5.2)
    add_bullets(
        slide,
        [
            "MCP server fronts the internal TiDB knowledge DB (docs + approved call transcripts).",
            "Retriever sends only relevant chunks and metadata to the model.",
            "Responses require citations back to internal chunk/document IDs.",
            "Policy filters and redaction run before external API calls.",
        ],
        7.15,
        2.15,
        5.2,
        3.8,
        14,
    )

    footer = slide.shapes.add_textbox(Inches(0.9), Inches(6.4), Inches(12.0), Inches(0.45))
    p = footer.text_frame.paragraphs[0]
    p.text = "Flow: user query -> policy checks -> MCP retrieval -> enterprise LLM -> cited output."
    p.font.name = "Aptos"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLORS["brand_dark"]
    p.alignment = PP_ALIGN.CENTER


def slide_3_outputs(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    style_background(slide)
    add_header(slide, "Guardrails for Prospect Calls", "How we protect sensitive prospect data with one LLM")

    add_panel(slide, 0.7, 1.4, 12.0, 5.0)
    add_bullets(
        slide,
        [
            "Only prospect calls with valid recording consent are eligible for ingestion and analysis.",
            "All-party consent states (for example CA and WA) must be handled with stricter notice controls.",
            "No existing customer private data is included unless it is already public marketing content.",
            "Internal-only access via OAuth/SSO; outbound messages restricted to @pingcap.com.",
            "Every generated asset includes retrieval citations and audit logs for traceability.",
        ],
        1.0,
        1.95,
        11.4,
        4.25,
        15,
    )


def slide_4_cost(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    style_background(slide)
    add_header(slide, "Cost Model (AWS VM + One Enterprise LLM)", "Simpler operating model with one model path")

    rows = 9
    cols = 3
    table = slide.shapes.add_table(rows, cols, Inches(0.95), Inches(1.6), Inches(11.7), Inches(4.45)).table
    table.columns[0].width = Inches(6.2)
    table.columns[1].width = Inches(2.4)
    table.columns[2].width = Inches(3.1)

    headers = ["Line Item", "Monthly (USD)", "Annual (USD)"]
    for i, value in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = value
        p = cell.text_frame.paragraphs[0]
        p.font.name = "Aptos"
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLORS["brand_dark"]

    data = [
        ("App/API/UI VM (t3.xlarge)", "$121.91", "$1,462.92"),
        ("Storage + networking + monitoring", "$79.72", "$956.64"),
        ("ChatGPT Business seats (20 users @ $25)", "$500.00", "$6,000.00"),
        ("Enterprise LLM API usage budget", "$250.00", "$3,000.00"),
        ("Backups + operating reserve", "$19.00", "$228.00"),
        ("Total operating cost", "$970.63", "$11,647.56"),
        ("With 20% contingency", "$1,164.76", "$13,977.07"),
        ("Typical single GTM technical hire", "-", "$200,000+"),
    ]

    for r, row in enumerate(data, start=1):
        for c, value in enumerate(row):
            cell = table.cell(r, c)
            cell.text = value
            p = cell.text_frame.paragraphs[0]
            p.font.name = "Aptos"
            p.font.size = Pt(12)
            p.font.color.rgb = COLORS["text"]
            if "Total" in row[0] or "contingency" in row[0]:
                p.font.bold = True
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(236, 242, 255)
            elif "Typical single GTM technical hire" in row[0]:
                p.font.bold = True
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(255, 245, 245)
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(255, 255, 255)

    note = slide.shapes.add_textbox(Inches(0.95), Inches(6.25), Inches(11.7), Inches(0.7))
    p = note.text_frame.paragraphs[0]
    p.text = "Seat baseline uses published ChatGPT Business pricing; swap this row with enterprise contract pricing when finalized."
    p.font.name = "Aptos"
    p.font.size = Pt(12)
    p.font.color.rgb = COLORS["muted"]


def slide_5_rollout_exit_sources(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    style_background(slide)
    add_header(slide, "Rollout + Exit Plan + Sources", "Start fast, measure value, and stop anytime if it is not working")

    add_panel(slide, 0.7, 1.4, 5.85, 4.95)
    panel_title(slide, "90-day rollout", 1.0, 1.78, 5.2)
    add_bullets(
        slide,
        [
            "Days 0-30: secure MVP, ingest core docs/transcripts, enable pilot reps and SEs.",
            "Days 31-60: launch coaching + asset workflows and tune retrieval quality.",
            "Days 61-90: scale to broader GTM teams with KPI reporting.",
        ],
        1.0,
        2.15,
        5.2,
        3.8,
        14,
    )

    add_panel(slide, 6.85, 1.4, 5.85, 4.95)
    panel_title(slide, "Cancel anytime + key sources", 7.15, 1.78, 5.2)
    add_bullets(
        slide,
        [
            "No long-term lock-in required for this architecture.",
            "If value is weak, we can stop quickly and retain internal artifacts/logs.",
            "OpenAI Enterprise Privacy (Jan 8, 2026): openai.com/enterprise-privacy",
            "OpenAI API data controls: platform.openai.com/docs/guides/your-data",
            "OpenAI Services Agreement (Jan 1, 2026): openai.com/policies/services-agreement",
            "MCP spec: modelcontextprotocol.io/introduction",
            "Consent laws: 18 U.S.C. 2511, CA Penal Code 632, WA RCW 9.73.030",
        ],
        7.15,
        2.15,
        5.2,
        3.8,
        14,
    )

    footer = slide.shapes.add_textbox(Inches(0.9), Inches(6.45), Inches(12.0), Inches(0.45))
    p = footer.text_frame.paragraphs[0]
    p.text = "Decision framing: high upside, low operating cost, and reversible at any point."
    p.font.name = "Aptos"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLORS["ok"]
    p.alignment = PP_ALIGN.CENTER


def build_deck() -> None:
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide_1_overview(prs)
    slide_2_single_llm_mcp(prs)
    slide_3_outputs(prs)
    slide_4_cost(prs)
    slide_5_rollout_exit_sources(prs)

    prs.save(OUTPUT_PATH)
    print(f"Deck generated at: {OUTPUT_PATH}")


if __name__ == "__main__":
    build_deck()
